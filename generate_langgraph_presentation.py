"""
LangGraph Agent Visualization - PowerPoint Generator

This script generates a Microsoft PowerPoint presentation (.pptx) that visualizes
and explains the LangGraph agent architectures used in the LLM Bootcamp Project.

Requirements:
    pip install python-pptx

Usage:
    python generate_langgraph_presentation.py

Output:
    langgraph_agents_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_title_slide(prs):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "LangGraph Agent Architecture"
    subtitle.text = "LLM Bootcamp Project\nAgentic RAG & ReAct Agents"

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)


def create_overview_slide(prs):
    """Create overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Add content
    p = tf.paragraphs[0]
    p.text = "Two Main LangGraph Implementations:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(20)

    items = [
        (
            "Agentic RAG Workflow",
            "Intelligent document question-answering with adaptive retrieval",
        ),
        (
            "ReAct Agent",
            "Web search-enabled conversational AI with tool-calling capability",
        ),
    ]

    for item_title, item_desc in items:
        p = tf.add_paragraph()
        p.text = f"• {item_title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.level = 0
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = item_desc
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(16)


def create_rag_architecture_slide(prs):
    """Create RAG architecture diagram slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Agentic RAG Workflow Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    # Node positions
    node_y = Inches(2)
    node_height = Inches(0.8)
    node_width = Inches(1.8)
    spacing = Inches(1.5)

    nodes = [
        ("START", Inches(0.8), RGBColor(100, 200, 100)),
        ("classify_mode", Inches(0.8) + spacing * 1, RGBColor(100, 150, 255)),
        ("retrieve", Inches(0.8) + spacing * 2, RGBColor(100, 150, 255)),
        ("generate", Inches(0.8) + spacing * 3, RGBColor(100, 150, 255)),
        ("END", Inches(0.8) + spacing * 4, RGBColor(255, 100, 100)),
    ]

    # Draw nodes
    prev_x = None
    for node_name, x_pos, color in nodes:
        # Create rounded rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, node_y, node_width, node_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(50, 50, 50)
        shape.line.width = Pt(2)

        # Add text
        text_frame = shape.text_frame
        text_frame.text = node_name
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Draw arrow from previous node
        if prev_x is not None:
            arrow_start_x = prev_x + node_width
            arrow_end_x = x_pos
            arrow_y = node_y + node_height / 2

            connector = slide.shapes.add_connector(
                1,  # Straight connector
                arrow_start_x,
                arrow_y,
                arrow_end_x,
                arrow_y,
            )
            connector.line.color.rgb = RGBColor(50, 50, 50)
            connector.line.width = Pt(2)

        prev_x = x_pos

    # Add descriptions below nodes
    descriptions = [
        ("Entry", Inches(0.8)),
        ("Classify query\ntype", Inches(0.8) + spacing * 1),
        ("Fetch docs\n(3 or 8)", Inches(0.8) + spacing * 2),
        ("Create\nresponse", Inches(0.8) + spacing * 3),
        ("Exit", Inches(0.8) + spacing * 4),
    ]

    desc_y = node_y + node_height + Inches(0.3)
    for desc_text, x_pos in descriptions:
        desc_box = slide.shapes.add_textbox(x_pos, desc_y, node_width, Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc_text
        desc_frame.paragraphs[0].font.size = Pt(11)
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_rag_nodes_slide(prs):
    """Create detailed RAG nodes explanation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "RAG Workflow Nodes"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    nodes_info = [
        ("classify_mode", "Analyzes query to determine 'summary' or 'fact' mode"),
        ("retrieve", "Fetches 8 docs (summary) or 3 docs (fact) from vector store"),
        ("generate", "Creates grounded response using only retrieved context"),
    ]

    for node_name, description in nodes_info:
        # Node name
        p = tf.paragraphs[0] if node_name == nodes_info[0][0] else tf.add_paragraph()
        p.text = f"🔹 {node_name}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(8)

        # Description
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(16)


def create_rag_state_slide(prs):
    """Create RAG state management slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "RAG State Management"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Add heading
    p = tf.paragraphs[0]
    p.text = "RAGState (TypedDict):"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(16)

    # State fields
    state_fields = [
        ("question: str", "User's query (input)"),
        ("mode: Literal['summary', 'fact']", "Response strategy"),
        ("documents: List[Document]", "Retrieved context"),
        ("generation: str", "Final response (output)"),
    ]

    for field, description in state_fields:
        p = tf.add_paragraph()
        p.text = field
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(204, 0, 0)
        p.level = 0
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f"→ {description}"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(12)


def create_rag_features_slide(prs):
    """Create RAG features slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "RAG Key Features"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    features = [
        ("Linear Workflow", "No conditional edges or loops - predictable flow"),
        ("Mode-Adaptive Retrieval", "Different strategies for summaries vs. facts"),
        ("Grounded Generation", "Responses strictly based on retrieved documents"),
        ("Type-Safe State", "TypedDict ensures proper data flow between nodes"),
        ("PII Anonymization", "Optional privacy protection using Presidio"),
        ("Vector Store Caching", "Faster reloads with FAISS cache (7-day TTL)"),
    ]

    p = tf.paragraphs[0]
    for i, (feature, description) in enumerate(features):
        if i > 0:
            p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(12)


def create_react_architecture_slide(prs):
    """Create ReAct agent architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "ReAct Agent Architecture (Cyclic)"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    # Create diagram with cycling capability
    center_x = Inches(5)
    center_y = Inches(3.5)

    # START node
    start_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        center_x - Inches(1),
        Inches(1.5),
        Inches(2),
        Inches(0.7),
    )
    start_shape.fill.solid()
    start_shape.fill.fore_color.rgb = RGBColor(100, 200, 100)
    start_shape.line.color.rgb = RGBColor(50, 50, 50)
    start_shape.text_frame.text = "START"
    start_shape.text_frame.paragraphs[0].font.size = Pt(16)
    start_shape.text_frame.paragraphs[0].font.bold = True
    start_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    start_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Agent Node
    agent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        center_x - Inches(1),
        center_y - Inches(0.5),
        Inches(2),
        Inches(1),
    )
    agent_shape.fill.solid()
    agent_shape.fill.fore_color.rgb = RGBColor(100, 150, 255)
    agent_shape.line.color.rgb = RGBColor(50, 50, 50)
    agent_shape.line.width = Pt(2.5)
    agent_shape.text_frame.text = "Agent Node\n(LLM Reasoning)"
    agent_shape.text_frame.paragraphs[0].font.size = Pt(16)
    agent_shape.text_frame.paragraphs[0].font.bold = True
    agent_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    agent_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Tools Node
    tools_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        center_x - Inches(1),
        center_y + Inches(1.5),
        Inches(2),
        Inches(1),
    )
    tools_shape.fill.solid()
    tools_shape.fill.fore_color.rgb = RGBColor(255, 180, 100)
    tools_shape.line.color.rgb = RGBColor(50, 50, 50)
    tools_shape.line.width = Pt(2.5)
    tools_shape.text_frame.text = "Tools Node\n(Tavily Search)"
    tools_shape.text_frame.paragraphs[0].font.size = Pt(16)
    tools_shape.text_frame.paragraphs[0].font.bold = True
    tools_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tools_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # END node
    end_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        center_x + Inches(2.5),
        center_y - Inches(0.35),
        Inches(1.5),
        Inches(0.7),
    )
    end_shape.fill.solid()
    end_shape.fill.fore_color.rgb = RGBColor(255, 100, 100)
    end_shape.line.color.rgb = RGBColor(50, 50, 50)
    end_shape.text_frame.text = "END"
    end_shape.text_frame.paragraphs[0].font.size = Pt(16)
    end_shape.text_frame.paragraphs[0].font.bold = True
    end_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    end_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Add text annotations
    annotations = [
        ("Tool call?", center_x + Inches(1.2), center_y + Inches(0.5), 12),
        ("Yes →", center_x - Inches(0.3), center_y + Inches(1.0), 12),
        ("← Results", center_x - Inches(1.5), center_y + Inches(1.0), 12),
        ("No →", center_x + Inches(1.2), center_y - Inches(0.1), 12),
    ]

    for text, x, y, size in annotations:
        text_box = slide.shapes.add_textbox(x, y, Inches(1), Inches(0.3))
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(size)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)


def create_react_features_slide(prs):
    """Create ReAct features slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ReAct Agent Features"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    features = [
        ("Cyclic Workflow", "Can loop multiple times for complex queries"),
        ("Tool Integration", "Tavily web search with 5 results per query"),
        ("Conditional Edges", "LLM decides when to use tools vs. respond"),
        ("Streaming Support", "Real-time response updates"),
        ("Timeout Protection", "90-second default timeout"),
        ("Multi-Step Reasoning", "Can chain multiple searches together"),
    ]

    p = tf.paragraphs[0]
    for i, (feature, description) in enumerate(features):
        if i > 0:
            p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(12)


def create_comparison_slide(prs):
    """Create comparison table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "RAG vs. ReAct Agent Comparison"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    # Create table
    rows = 8
    cols = 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.25)
    table.columns[2].width = Inches(3.25)

    # Header row
    headers = ["Feature", "Agentic RAG", "ReAct Agent"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    data = [
        ["Purpose", "Document QA", "Web search + Chat"],
        ["Data Source", "Uploaded PDFs", "Real-time web (Tavily)"],
        ["Graph Type", "Linear (fixed)", "Cyclic (loops)"],
        ["Nodes", "3 agent nodes", "2 nodes (agent + tools)"],
        ["Conditional Logic", "None", "Yes (tool decisions)"],
        ["Response Mode", "Adaptive (2 modes)", "Single strategy"],
        ["Streaming", "Not implemented", "Supported"],
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(11)

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)


def create_use_cases_slide(prs):
    """Create use cases slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "When to Use Each Agent"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # RAG use cases
    p = tf.paragraphs[0]
    p.text = "🟦 Agentic RAG - Use When:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_after = Pt(10)

    rag_cases = [
        "You have proprietary documents (PDFs, reports)",
        "Need factually grounded responses",
        "Require citation to source material",
        "Want consistent, reproducible answers",
    ]

    for case in rag_cases:
        p = tf.add_paragraph()
        p.text = f"• {case}"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(8)

    # ReAct use cases
    p = tf.add_paragraph()
    p.text = "🟧 ReAct Agent - Use When:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)
    p.space_after = Pt(10)
    p.space_before = Pt(20)

    react_cases = [
        "Need current information (news, prices, events)",
        "Want broad web knowledge",
        "Require multi-step reasoning",
        "Need to combine multiple sources",
    ]

    for case in react_cases:
        p = tf.add_paragraph()
        p.text = f"• {case}"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(8)


def create_code_example_slide(prs):
    """Create code example slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Implementation Example"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Code example
    code = """# RAG Workflow Setup
from langchain_helpers import RAGHelper

rag_workflow, pii_entities = RAGHelper.setup_rag_system(
    uploaded_files,
    api_key=openai_key,
    anonymize_pii=True,
    use_cache=True
)

# Process query
result = rag_workflow.invoke({"question": user_query})
answer = result["generation"]

# ReAct Agent Setup
from langchain_helpers import AgentChatbotHelper

agent = AgentChatbotHelper.setup_agent(
    openai_api_key, tavily_api_key
)
response = await AgentChatbotHelper.process_agent_response(
    agent, user_query
)"""

    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.space_after = Pt(10)


def create_file_references_slide(prs):
    """Create file references slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Implementation File References"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    references = [
        (
            "langchain_helpers.py",
            [
                "Lines 591-728: RAGHelper.build_simple_agentic_rag()",
                "Lines 730-781: RAGHelper.setup_rag_system()",
                "Lines 129-228: AgentChatbotHelper (ReAct)",
            ],
        ),
        (
            "Streamlit Pages",
            ["pages/3_📄_RAG_Document_Chat.py", "pages/2_🔎_Search-Enabled_Chat.py"],
        ),
        (
            "Configuration",
            [
                "config.py: Secure API key management",
                "ui_components.py: Reusable UI components",
            ],
        ),
    ]

    p = tf.paragraphs[0]
    for i, (section, items) in enumerate(references):
        if i > 0:
            p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(8)
        p.space_before = Pt(12) if i > 0 else Pt(0)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.name = "Courier New"
            p.level = 1
            p.space_after = Pt(6)


def create_summary_slide(prs):
    """Create summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Takeaways"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    takeaways = [
        "Both agents use typed state management for reliability",
        "RAG provides grounded, factual responses from documents",
        "ReAct enables real-time web search and multi-step reasoning",
        "Graphs are compiled and optimized by LangGraph",
        "Production-ready with caching, PII protection, and error handling",
        "Easily extensible with custom nodes and tools",
    ]

    p = tf.paragraphs[0]
    for i, takeaway in enumerate(takeaways):
        if i > 0:
            p = tf.add_paragraph()
        p.text = f"✓ {takeaway}"
        p.font.size = Pt(18)
        p.space_after = Pt(14)


def create_next_steps_slide(prs):
    """Create next steps slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    steps = [
        ("Try the Agents", "Run: streamlit run Home.py"),
        ("Upload Documents", "Test RAG workflow with your PDFs"),
        ("Web Search", "Try ReAct agent for current events"),
        ("Customize", "Modify node logic in langchain_helpers.py"),
        ("Explore Notebook", "See langgraph_visualization.ipynb"),
        ("Read Docs", "Check CLAUDE.md for full architecture"),
    ]

    p = tf.paragraphs[0]
    for i, (step, description) in enumerate(steps):
        if i > 0:
            p = tf.add_paragraph()
        p.text = f"{i + 1}. {step}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(15)
        p.level = 1
        p.space_after = Pt(12)


def main():
    """Generate the PowerPoint presentation."""
    print("🎨 Generating LangGraph Agents PowerPoint Presentation...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    print("  📄 Creating title slide...")
    create_title_slide(prs)

    print("  📄 Creating overview slide...")
    create_overview_slide(prs)

    print("  📄 Creating RAG architecture diagram...")
    create_rag_architecture_slide(prs)

    print("  📄 Creating RAG nodes explanation...")
    create_rag_nodes_slide(prs)

    print("  📄 Creating RAG state management slide...")
    create_rag_state_slide(prs)

    print("  📄 Creating RAG features slide...")
    create_rag_features_slide(prs)

    print("  📄 Creating ReAct architecture diagram...")
    create_react_architecture_slide(prs)

    print("  📄 Creating ReAct features slide...")
    create_react_features_slide(prs)

    print("  📄 Creating comparison table...")
    create_comparison_slide(prs)

    print("  📄 Creating use cases slide...")
    create_use_cases_slide(prs)

    print("  📄 Creating code example slide...")
    create_code_example_slide(prs)

    print("  📄 Creating file references slide...")
    create_file_references_slide(prs)

    print("  📄 Creating summary slide...")
    create_summary_slide(prs)

    print("  📄 Creating next steps slide...")
    create_next_steps_slide(prs)

    # Save presentation
    output_file = "langgraph_agents_presentation.pptx"
    prs.save(output_file)

    print("\n✅ Presentation created successfully!")
    print(f"📁 Output file: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(
        "\nTo view the presentation, open it in Microsoft PowerPoint or compatible software."
    )


if __name__ == "__main__":
    main()
